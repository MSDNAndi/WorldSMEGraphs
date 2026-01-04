#!/usr/bin/env python3
"""
PowerPoint and PDF Presentation Generator for WorldSMEGraphs
=============================================================

Creates actual .pptx PowerPoint and .pdf files from presentation content
with AI-generated images for each slide.

Features:
- Creates real PowerPoint (.pptx) files using python-pptx
- Converts to PDF using reportlab
- Integrates with GPT Image 1.5 for slide visuals
- British humor style support (sarcastic, self-deprecating, quirky)
- Content safety handling with prompt modification

Dependencies:
- python-pptx: PowerPoint file generation
- reportlab: PDF generation
- PIL/Pillow: Image handling

Usage:
    python presentation_generator.py --slides slides.yaml --output presentation
    python presentation_generator.py --slides slides.yaml --style british-humor
    python presentation_generator.py --slides slides.yaml --generate-images

Author: WorldSMEGraphs Presentation Agent
Version: 1.0.0
Created: 2026-01-04
"""

import asyncio
import argparse
from pathlib import Path
from typing import List, Optional, Tuple
from dataclasses import dataclass, field

# PowerPoint library
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    RGBColor = None
    print("Warning: python-pptx not installed. Run: pip install python-pptx")

# PDF library  
try:
    from reportlab.lib.pagesizes import LETTER, landscape
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.lib.colors import HexColor
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    print("Warning: reportlab not installed. Run: pip install reportlab")

# Image handling (PIL availability check)
try:
    import PIL
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: Pillow not installed. Run: pip install Pillow")

# YAML for slide definitions
try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False

# =============================================================================
# CONSTANTS
# =============================================================================

# British Humor Style Guide
BRITISH_HUMOR_GUIDELINES = """
British Humor Style for Technical Presentations
================================================

KEY CHARACTERISTICS:
1. SELF-DEPRECATING: "Well, I've only been getting this wrong for 15 years..."
2. UNDERSTATED: "This is slightly important" (for critical concepts)
3. DRY SARCASM: "Oh good, another abstraction. Just what we needed."
4. IRONIC OBSERVATIONS: "Monads: solving problems you didn't know you had"
5. POLITE CRITICISM: "With the greatest respect, this is absolute rubbish"
6. QUINTESSENTIALLY BRITISH REFERENCES: Tea, queues, weather, apologizing

TONE MARKERS:
- "Rather" / "Quite" / "Slightly" (understatement)
- "One might argue..." (polite disagreement)
- "As it were" / "So to speak" (hedging)
- "Right then" / "Well then" (transitions)
- "Brilliant" (often sarcastic)
- "Sorry, but..." (apologetic directness)

VISUAL HUMOR (PG-friendly):
- Confused stick figures with tea cups
- Queue diagrams (very British)
- Weather metaphors (unpredictable like async code)
- Garden shed analogies (British inventions)
- The Tube map style for concept relationships

WHAT TO AVOID:
- Crude or offensive content
- Punching down
- Excessive cynicism (British humor is warm underneath)
- Overused stereotypes
- American humor styles (too direct)
"""

# Content safety retry prompts
CONTENT_SAFETY_ALTERNATIVES = {
    "violence": "Replace action with gentle transformation or peaceful transition",
    "adult": "Use abstract shapes and professional imagery instead",
    "controversial": "Focus on neutral technical concepts and universal themes",
    "political": "Remove any political references, use generic symbols",
    "medical": "Use abstract anatomical representations without specific organs",
}

# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class SlideContent:
    """Content for a single slide."""
    title: str
    body: str = ""
    speaker_notes: str = ""
    image_prompt: str = ""
    image_path: Optional[Path] = None
    humor_style: str = "british"
    layout: str = "title_and_content"  # title_only, title_and_content, two_column, image_full

@dataclass
class PresentationConfig:
    """Configuration for presentation generation."""
    title: str
    subtitle: str = ""
    author: str = "WorldSMEGraphs"
    slides: List[SlideContent] = field(default_factory=list)
    output_path: Path = field(default_factory=lambda: Path("./presentation"))
    generate_images: bool = True
    style: str = "british-humor"  # british-humor, professional, fun
    theme_color: str = "#0078D4"  # Azure blue by default
    image_dir: Path = field(default_factory=lambda: Path("./generated_images"))


# =============================================================================
# CONTENT SAFETY HANDLER (Import from gpt_image_generator for shared code)
# =============================================================================

# Try to import ContentSafetyHandler from gpt_image_generator
try:
    from gpt_image_generator import ContentSafetyHandler
    CONTENT_SAFETY_AVAILABLE = True
except ImportError:
    # Fallback: Define minimal version if import fails
    CONTENT_SAFETY_AVAILABLE = False
    import re
    
    class ContentSafetyHandler:
        """Fallback ContentSafetyHandler when gpt_image_generator is not available."""
        
        SAFE_REPLACEMENTS = {
            "death": "transition",
            "kill": "terminate gracefully",
            "blood": "energy flow",
            "attack": "challenge",
            "crash": "unexpected stop",
            "destroy": "clean up",
            "explode": "expand rapidly",
            "injection": "insertion",
            "execute": "run",
            "daemon": "background service",
            "master": "primary",
            "slave": "replica",
        }
        
        @classmethod
        def sanitize_prompt(cls, prompt: str) -> str:
            """Pre-emptively sanitize a prompt before sending to the API."""
            result = prompt
            for old, new in cls.SAFE_REPLACEMENTS.items():
                if old.lower() in result.lower():
                    pattern = re.compile(re.escape(old), re.IGNORECASE)
                    result = pattern.sub(new, result)
            return result
        
        @classmethod
        def handle_safety_error(cls, original_prompt: str, error_message: str) -> Tuple[str, bool]:
            """Handle a content safety error by modifying the prompt."""
            modified = cls.sanitize_prompt(original_prompt)
            safety_suffix = """

Style: Clean, professional, abstract representation.
Constraints:
- Family-friendly content only
- Professional corporate style
"""
            return modified + safety_suffix, True


# =============================================================================
# POWERPOINT GENERATOR
# =============================================================================

class PowerPointGenerator:
    """
    Generates PowerPoint (.pptx) files from slide content.
    """
    
    def __init__(self, config: PresentationConfig):
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required. Run: pip install python-pptx")
        self.config = config
        self.prs = Presentation()
        self._setup_slide_dimensions()
        
    def _setup_slide_dimensions(self):
        """Set slide dimensions to 16:9 widescreen."""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def _get_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGB."""
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def add_title_slide(self):
        """Add the title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = self.config.title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        if self.config.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = self.config.subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.alignment = PP_ALIGN.CENTER
            
        # Add author
        author_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5)
        )
        author_frame = author_box.text_frame
        author_para = author_frame.paragraphs[0]
        author_para.text = self.config.author
        author_para.font.size = Pt(14)
        author_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(self, content: SlideContent):
        """Add a content slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Layout calculations
        if content.image_path and Path(content.image_path).exists():
            # Two-column layout with image
            text_width = Inches(6)
            text_left = Inches(0.5)
            image_left = Inches(7)
            image_width = Inches(5.833)
        else:
            # Full-width text
            text_width = Inches(12.333)
            text_left = Inches(0.5)
            image_left = None
            image_width = None
        
        # Add title
        title_box = slide.shapes.add_textbox(
            text_left, Inches(0.5), text_width, Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        # Add body content
        if content.body:
            body_box = slide.shapes.add_textbox(
                text_left, Inches(1.7), text_width, Inches(5)
            )
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            
            # Split body into bullet points
            lines = content.body.strip().split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    para = body_frame.paragraphs[0]
                else:
                    para = body_frame.add_paragraph()
                
                para.text = line.lstrip("•-* ")
                para.font.size = Pt(18)
                para.level = 0 if line.strip().startswith(("•", "-", "*")) else 0
        
        # Add image if available
        if image_left and content.image_path and Path(content.image_path).exists():
            try:
                slide.shapes.add_picture(
                    str(content.image_path),
                    image_left, Inches(1.5),
                    width=image_width
                )
            except Exception as e:
                print(f"Warning: Could not add image {content.image_path}: {e}")
        
        # Add speaker notes
        if content.speaker_notes:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            notes_frame.text = content.speaker_notes
        
        return slide
    
    def generate(self) -> Path:
        """Generate the PowerPoint file."""
        # Add title slide
        self.add_title_slide()
        
        # Add content slides
        for slide_content in self.config.slides:
            self.add_content_slide(slide_content)
        
        # Save
        output_path = self.config.output_path.with_suffix(".pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        
        print(f"✓ PowerPoint saved: {output_path}")
        return output_path


# =============================================================================
# PDF GENERATOR
# =============================================================================

class PDFGenerator:
    """
    Generates PDF files from slide content using ReportLab.
    """
    
    def __init__(self, config: PresentationConfig):
        if not REPORTLAB_AVAILABLE:
            raise ImportError("reportlab is required. Run: pip install reportlab")
        self.config = config
        self._setup_styles()
        
    def _setup_styles(self):
        """Setup PDF styles."""
        self.styles = getSampleStyleSheet()
        
        # Title style
        self.styles.add(ParagraphStyle(
            name='SlideTitle',
            parent=self.styles['Heading1'],
            fontSize=28,
            alignment=TA_CENTER,
            spaceAfter=30,
            textColor=HexColor(self.config.theme_color),
        ))
        
        # Body style
        self.styles.add(ParagraphStyle(
            name='SlideBody',
            parent=self.styles['Normal'],
            fontSize=14,
            alignment=TA_LEFT,
            spaceBefore=10,
            spaceAfter=10,
            leftIndent=20,
        ))
        
        # Subtitle style
        self.styles.add(ParagraphStyle(
            name='SlideSubtitle',
            parent=self.styles['Normal'],
            fontSize=18,
            alignment=TA_CENTER,
            spaceBefore=20,
        ))
        
    def generate(self) -> Path:
        """Generate the PDF file."""
        output_path = self.config.output_path.with_suffix(".pdf")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=landscape(LETTER),
            rightMargin=0.5*inch,
            leftMargin=0.5*inch,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch,
        )
        
        story = []
        
        # Title page
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph(self.config.title, self.styles['SlideTitle']))
        if self.config.subtitle:
            story.append(Paragraph(self.config.subtitle, self.styles['SlideSubtitle']))
        story.append(Spacer(1, inch))
        story.append(Paragraph(self.config.author, self.styles['SlideSubtitle']))
        story.append(PageBreak())
        
        # Content slides
        for slide in self.config.slides:
            # Title
            story.append(Paragraph(slide.title, self.styles['SlideTitle']))
            
            # Image if available
            if slide.image_path and Path(slide.image_path).exists():
                try:
                    img = RLImage(str(slide.image_path), width=4*inch, height=3*inch)
                    story.append(img)
                    story.append(Spacer(1, 0.25*inch))
                except Exception as e:
                    print(f"Warning: Could not add image to PDF: {e}")
            
            # Body content
            if slide.body:
                for line in slide.body.strip().split("\n"):
                    clean_line = line.lstrip("•-* ")
                    if clean_line:
                        story.append(Paragraph(f"• {clean_line}", self.styles['SlideBody']))
            
            story.append(PageBreak())
        
        doc.build(story)
        print(f"✓ PDF saved: {output_path}")
        return output_path


# =============================================================================
# BRITISH HUMOR CONTENT GENERATOR
# =============================================================================

def generate_british_humor_slides(topic: str = "Functional Programming") -> List[SlideContent]:
    """
    Generate slide content with British humor style.
    
    This is specifically for the Functional Programming presentation
    with a sarcastic, self-deprecating, quirky British tone.
    """
    slides = [
        # Slide 1: Title
        SlideContent(
            title="Functional Programming: A Terribly Civilised Introduction",
            body="Or: How I Learned to Stop Worrying and Love the Monad",
            speaker_notes="Open with a self-deprecating joke about how long it took you to understand this.",
            image_prompt="""Create a professional presentation title slide background.

Scene:
A stylized British tea cup (porcelain, floral pattern) sitting on a saucer, with 
abstract code symbols (lambdas λ, arrows →) floating gently in the steam rising 
from the cup. The cup is positioned in the lower right corner.

Style:
- Clean, minimal, modern digital illustration
- Soft pastel colors with British racing green accents (#004225)
- Gentle, calming atmosphere
- Professional corporate presentation quality

Technical:
- 16:9 aspect ratio composition
- Leave left side mostly empty for text overlay
- Subtle geometric patterns in background

Constraints:
- Family-friendly, professional
- No text in the image
- Original artwork only
""",
            layout="title_only"
        ),
        
        # Slide 2: Honest Opener
        SlideContent(
            title="Right Then, Let's Be Honest",
            body="""• You're here because someone said "monads" in a meeting
• And you nodded wisely whilst dying inside
• Don't worry, we've all been there
• (Some of us for rather longer than we'd care to admit)
• Today we shall fix that. Probably. Maybe. We'll see.""",
            speaker_notes="This is where you establish rapport. Everyone is confused. That's normal.",
            image_prompt="""Create a humorous illustration for a technical presentation.

Scene:
A simple stick figure character sitting at a desk, looking confused at a 
computer screen. Above their head is a thought bubble containing a question 
mark. On the desk is a cup of tea (British detail). The character's 
expression is bewildered but friendly.

Style:
- Simple, clean line art illustration
- Minimalist style with touches of humor
- Soft blue and grey tones
- Professional but approachable

Technical:
- Centered composition
- Clean white background
- Character should look relatable, not silly

Constraints:
- Family-friendly
- No text in image
- Original artwork, not clip art
"""
        ),
        
        # Slide 3: Category Theory (Gentle Introduction)
        SlideContent(
            title="Category Theory: It's Less Scary Than It Sounds",
            body="""• Invented by mathematicians in the 1940s
• (They needed something to do during the war besides codebreaking)
• It's really just about "things" and "arrows between things"
• Objects + Morphisms = That's literally it
• You've been doing this all along, you just didn't know it had a fancy name""",
            speaker_notes="Demystify immediately. Most developers already understand the core concepts.",
            image_prompt="""Create an educational diagram for category theory basics.

Scene:
Three simple geometric shapes (circle, square, triangle) arranged in a 
triangle formation. Connecting them are curved arrows pointing LEFT TO RIGHT 
between each pair. Each arrow is labeled with a simple symbol (f, g, h).
The overall feeling is friendly and approachable.

Style:
- Clean, modern educational illustration
- Soft pastel colors: circle=soft blue, square=soft green, triangle=soft coral
- Arrows are smooth curves with clear arrowheads on the RIGHT side
- Gentle drop shadows for depth

Technical:
- Arrows show composition direction: f: A→B means arrow FROM A TO B
- Keep layout balanced and symmetrical
- Professional presentation quality

Constraints:
- No text labels (just symbolic arrows)
- Family-friendly, educational
- Original artwork only
"""
        ),
        
        # Slide 4: Functors
        SlideContent(
            title="Functors: The Art of Preserving Structure",
            body="""• A Functor is a "structure-preserving mapping"
• (In plainer English: .Select() in LINQ, .map() in JavaScript)
• It takes things in boxes and transforms them whilst keeping them boxed
• The box might be: List, Option, Task, Observable...
• Quite useful, really. Terribly sensible.""",
            speaker_notes="Connect to everyday code. Everyone uses map. Functors are just the formal name.",
            image_prompt="""Create a diagram showing functor concept for developers.

Scene:
Two rows of boxes (containers). Top row: three boxes containing simple shapes
(star, heart, circle). Bottom row: the same three boxes, but the shapes inside 
have been transformed (star→pentagon, heart→diamond, circle→hexagon).
Between the rows, arrows pointing DOWNWARD show the transformation.
The boxes themselves remain identical - only contents change.

Style:
- Clean, modern diagram style
- Soft blue boxes with white interiors
- Transformation arrows are dashed and purple
- Gentle, professional aesthetic

Technical:
- Arrows point DOWN to show transformation direction
- Boxes are identical top and bottom (structure preserved)
- Transformation is clearly visible

Constraints:
- No text
- Family-friendly
- Original artwork only
"""
        ),
        
        # Slide 5: Monoids
        SlideContent(
            title="Monoids: Things That Combine Nicely",
            body="""• Two rules: things can be combined, and there's an "identity" thing
• Numbers with addition: 1 + 2 = 3, and 0 + anything = anything
• Strings with concatenation: "Hello" + "World", and "" + anything = anything
• Lists: [1,2] ++ [3,4], and [] ++ anything = anything
• It's like the Queen's rules of etiquette, but for data""",
            speaker_notes="Monoids are everywhere. Make it click with familiar examples.",
            image_prompt="""Create an illustration showing monoid concept (combining things).

Scene:
A visual metaphor of proper British tea service. Three tea cups are shown
merging together (LEFT TO RIGHT flow) into one elegant teapot. Next to the
teapot is an empty cup labeled implicitly as "identity" (adds nothing).
The flow shows combination while maintaining the "tea" nature.

Style:
- Elegant, British aesthetic with a touch of whimsy
- Soft cream and china blue colors
- Flowing arrows showing combination direction (LEFT TO RIGHT)
- Clean, minimalist illustration

Technical:
- Flow direction is explicitly LEFT TO RIGHT
- Empty cup (identity) is visually distinct but same style
- Combination is harmonious, not chaotic

Constraints:
- No text in image
- Family-friendly
- Original artwork only
"""
        ),
        
        # Slide 6: Monads - The Big One
        SlideContent(
            title="Monads: Sorry About This One",
            body="""• Right, deep breath. Here we go.
• A Monad is a Functor that can also "flatten" nested structures
• .SelectMany() in LINQ, .flatMap() in JavaScript
• Promise.then() chains? That's a Monad.
• Optional?.value chaining? Also a Monad.
• You've been using Monads for years. Surprise!""",
            speaker_notes="The big reveal. They already know this. Just didn't know the name.",
            image_prompt="""Create an illustration for monad concept (nested containers that flatten).

Scene:
A visual progression showing three stages from LEFT TO RIGHT:
1. LEFT: A box inside a box inside a box (nested, like Russian dolls)
2. MIDDLE: An arrow labeled with transformation, pointing RIGHT
3. RIGHT: A single box containing the inner content (flattened)
The progression clearly shows unwrapping/flattening while preserving content.

Style:
- Clean, modern diagram style
- Soft gradient blues for the boxes (darker=outer, lighter=inner)
- Smooth transformation arrows pointing RIGHT
- Professional presentation quality

Technical:
- Flow is explicitly LEFT TO RIGHT
- Nesting is clearly visible on the left
- Flattening result is clean on the right
- The content inside remains unchanged

Constraints:
- No text
- Family-friendly
- Original artwork only
"""
        ),
        
        # Slide 7: The Grand Unification
        SlideContent(
            title="Putting It All Together (Without Panic)",
            body="""• Category Theory → The framework for thinking about structure
• Functors → Transform contents, preserve structure  
• Monoids → Combine things sensibly
• Monads → Functors that handle nested complexity
• Together: A complete vocabulary for composition
• Like LEGO, but for your code. Rather satisfying, actually.""",
            speaker_notes="Tie it all together. This is the 'aha' moment.",
            image_prompt="""Create a grand unification diagram for functional programming concepts.

Scene:
Four distinct elements (representing Category Theory, Functors, Monoids, Monads)
arranged in a circle, with elegant curved lines connecting them all to a 
central glowing point. Each element has a distinct shape and soft pastel color.
The overall composition suggests harmony and interconnection.

Style:
- Elegant, sophisticated diagram
- Soft pastel colors: blue, green, coral, purple
- Golden central connection point
- Clean, modern, professional

Technical:
- Circular arrangement suggests equality of importance
- Lines flow smoothly, no sharp angles
- Central point is clearly the unifying element

Constraints:
- No text
- Family-friendly
- Original artwork only
"""
        ),
        
        # Slide 8: Practical Takeaways
        SlideContent(
            title="What to Actually Do With This",
            body="""• Use .map() and .flatMap() with confidence now
• Recognize patterns: "Ah, that's a Functor situation"
• Compose small functions into larger pipelines
• Handle errors with Either/Result types (Monads!)
• Impress colleagues with restrained accuracy
• Make better tea. (That last one is optional but recommended.)""",
            speaker_notes="Actionable items. Send them off with something practical.",
            image_prompt="""Create an illustration showing practical developer victory.

Scene:
A simple, friendly stick figure character at a desk, now looking confident
and composed. They're making a "thumbs up" gesture. On the screen is a 
simple code pattern. A cup of tea sits completed on the desk. 
Above their head is a small lightbulb (idea understood).

Style:
- Clean, minimal line art
- Warm, encouraging colors
- Professional but friendly
- Celebratory without being over the top

Technical:
- Character centered
- Simple, clean composition
- Warm yellow/gold accents for success

Constraints:
- Family-friendly
- No text in image
- Original artwork only
"""
        ),
        
        # Slide 9: Q&A
        SlideContent(
            title="Any Questions?",
            body="""(Besides "Why did this take so long to explain?")

• Perfectly reasonable to still feel slightly confused
• Rome wasn't built in a day, and neither was Haskell comprehension
• Resources for further reading in the appendix
• Feel free to pretend you understood everything. We won't tell.""",
            speaker_notes="Open for Q&A. Have backup jokes ready for awkward silences.",
            image_prompt="""Create an illustration for Q&A slide.

Scene:
A friendly question mark character (stylized, with a face and arms) standing
at a podium, looking welcoming. Around it are raised hands (abstract, 
friendly shapes) suggesting an engaged audience. The overall mood is 
inviting and non-intimidating.

Style:
- Warm, friendly illustration style
- Soft blues and gentle yellows
- Approachable, not corporate
- Slightly whimsical

Technical:
- Question mark is the central character
- Hands are abstract, not realistic
- Composition is balanced and inviting

Constraints:
- Family-friendly
- No text in image  
- Original artwork only
"""
        ),
        
        # Slide 10: Thank You
        SlideContent(
            title="Thank You for Your Patience",
            body="""You've been a wonderful audience.

Now go forth and map things!

[Tea service illustration]""",
            speaker_notes="Wrap up warmly. Thank the audience genuinely.",
            image_prompt="""Create a closing slide illustration with British warmth.

Scene:
An elegant, stylized British teapot pouring tea into a cup, with gentle
steam rising. Around it are abstract symbols of code (λ, →, ∘) floating
peacefully like leaves. The scene conveys completion and satisfaction.

Style:
- Warm, elegant illustration
- Soft cream and British racing green colors
- Peaceful, satisfying atmosphere
- Professional presentation quality

Technical:
- Centered, balanced composition
- Warm lighting effect
- Symbols are subtle, not overwhelming

Constraints:
- Family-friendly
- No text in image
- Original artwork only
"""
        ),
    ]
    
    return slides


# =============================================================================
# MAIN GENERATOR
# =============================================================================

async def generate_presentation(
    config: PresentationConfig,
    generate_images: bool = True,
    image_generator_func = None
) -> Tuple[Optional[Path], Optional[Path]]:
    """
    Generate a complete presentation with optional images.
    
    Args:
        config: Presentation configuration
        generate_images: Whether to generate images for slides
        image_generator_func: Async function to generate images
        
    Returns:
        Tuple of (pptx_path, pdf_path)
    """
    print(f"\n🎨 Generating presentation: {config.title}")
    print(f"   Style: {config.style}")
    print(f"   Slides: {len(config.slides)}")
    
    # Generate images if requested and generator available
    if generate_images and image_generator_func:
        print("\n📸 Generating slide images...")
        for i, slide in enumerate(config.slides):
            if slide.image_prompt:
                print(f"   [{i+1}/{len(config.slides)}] {slide.title[:40]}...")
                try:
                    # Apply content safety
                    safe_prompt = ContentSafetyHandler.sanitize_prompt(slide.image_prompt)
                    
                    # Generate image
                    image_path = await image_generator_func(
                        safe_prompt,
                        output_prefix=f"slide_{i+1:02d}",
                        output_dir=config.image_dir
                    )
                    slide.image_path = image_path
                except Exception as e:
                    print(f"   ⚠️  Image generation failed: {e}")
                    # Try with safer prompt
                    try:
                        safer_prompt, _ = ContentSafetyHandler.handle_safety_error(
                            slide.image_prompt, str(e)
                        )
                        image_path = await image_generator_func(
                            safer_prompt,
                            output_prefix=f"slide_{i+1:02d}_retry",
                            output_dir=config.image_dir
                        )
                        slide.image_path = image_path
                    except Exception as e2:
                        print(f"   ❌ Retry also failed: {e2}")
    
    pptx_path = None
    pdf_path = None
    
    # Generate PowerPoint
    if PPTX_AVAILABLE:
        try:
            pptx_gen = PowerPointGenerator(config)
            pptx_path = pptx_gen.generate()
        except Exception as e:
            print(f"⚠️  PowerPoint generation failed: {e}")
    
    # Generate PDF
    if REPORTLAB_AVAILABLE:
        try:
            pdf_gen = PDFGenerator(config)
            pdf_path = pdf_gen.generate()
        except Exception as e:
            print(f"⚠️  PDF generation failed: {e}")
    
    return pptx_path, pdf_path


# =============================================================================
# CLI INTERFACE
# =============================================================================

def main():
    """CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint and PDF presentations with AI images",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=BRITISH_HUMOR_GUIDELINES
    )
    
    parser.add_argument("--title", "-t", default="Functional Programming: A British Introduction",
                       help="Presentation title")
    parser.add_argument("--subtitle", default="A Terribly Civilised Guide to Category Theory",
                       help="Presentation subtitle")
    parser.add_argument("--output", "-o", type=Path, default=Path("./presentation"),
                       help="Output path (without extension)")
    parser.add_argument("--style", choices=["british-humor", "professional", "fun"],
                       default="british-humor", help="Presentation style")
    parser.add_argument("--generate-images", "-g", action="store_true",
                       help="Generate images for slides (requires image generator)")
    parser.add_argument("--image-dir", type=Path, default=Path("./generated_images"),
                       help="Directory for generated images")
    parser.add_argument("--slides-file", type=Path,
                       help="YAML file with slide definitions")
    
    args = parser.parse_args()
    
    # Build configuration
    if args.slides_file and args.slides_file.exists() and YAML_AVAILABLE:
        with open(args.slides_file) as f:
            slides_data = yaml.safe_load(f)
        slides = [SlideContent(**s) for s in slides_data.get("slides", [])]
    else:
        # Use default British humor slides
        slides = generate_british_humor_slides()
    
    config = PresentationConfig(
        title=args.title,
        subtitle=args.subtitle,
        slides=slides,
        output_path=args.output,
        style=args.style,
        image_dir=args.image_dir,
        generate_images=args.generate_images,
    )
    
    # Run generation (without images for now - images need the image generator)
    pptx_path, pdf_path = asyncio.run(generate_presentation(
        config, 
        generate_images=False  # Images handled separately
    ))
    
    print("\n✅ Generation complete!")
    if pptx_path:
        print(f"   PowerPoint: {pptx_path}")
    if pdf_path:
        print(f"   PDF: {pdf_path}")


if __name__ == "__main__":
    main()
