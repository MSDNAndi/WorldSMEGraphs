#!/usr/bin/env python3
"""
Simple presentation builder for Category Theory presentation
Creates PPTX and PDF from pre-generated images
"""

from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape, letter
from PIL import Image
import os

def create_presentation():
    """Create PowerPoint presentation with pre-generated images"""
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Get image files
    image_dir = "images"
    images = sorted([f for f in os.listdir(image_dir) if f.endswith('.png')])
    
    print(f"Found {len(images)} images:")
    for img in images:
        print(f"  - {img}")
    
    # Add slides with images
    for image_file in images:
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add image to fill the slide
        image_path = os.path.join(image_dir, image_file)
        left = Inches(0)
        top = Inches(0)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(16))
        
        print(f"Added slide: {image_file}")
    
    # Save PPTX
    pptx_file = "professional-category-theory.pptx"
    prs.save(pptx_file)
    print(f"\n✓ Created: {pptx_file}")
    
    return pptx_file, images

def create_pdf(images):
    """Create PDF from images"""
    pdf_file = "professional-category-theory.pdf"
    
    # Create PDF with landscape letter size (11x8.5 inches)
    c = canvas.Canvas(pdf_file, pagesize=landscape(letter))
    width, height = landscape(letter)
    
    for image_file in images:
        image_path = os.path.join("images", image_file)
        
        # Open image to get dimensions
        img = Image.open(image_path)
        img_width, img_height = img.size
        
        # Calculate scaling to fit page while maintaining aspect ratio
        scale = min(width / img_width, height / img_height)
        scaled_width = img_width * scale
        scaled_height = img_height * scale
        
        # Center image on page
        x = (width - scaled_width) / 2
        y = (height - scaled_height) / 2
        
        # Draw image
        c.drawImage(image_path, x, y, width=scaled_width, height=scaled_height)
        c.showPage()
        
        print(f"Added PDF page: {image_file}")
    
    c.save()
    print(f"\n✓ Created: {pdf_file}")
    
    return pdf_file

if __name__ == "__main__":
    print("Building Professional Category Theory Presentation")
    print("=" * 60)
    
    # Create PPTX
    pptx_file, images = create_presentation()
    
    # Create PDF
    pdf_file = create_pdf(images)
    
    print("\n" + "=" * 60)
    print("✓ Presentation build complete!")
    print(f"\nGenerated files:")
    print(f"  - {pptx_file}")
    print(f"  - {pdf_file}")
    print("\nSlides:")
    for i, img in enumerate(images, 1):
        print(f"  {i}. {img}")
